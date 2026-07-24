#!/usr/bin/env python3
"""
读取参考资料文件，支持多种格式，统一输出为文本 + 图片索引。

设计原则：
- 文字提取：用 markitdown / python-docx / PyPDF2 / python-pptx
- 图片提取：用 python-docx / PyMuPDF / python-pptx 导出嵌入图片
- 图片理解：不做 OCR，交给多模态大模型用 Read 工具看导出的图片
- 装饰性图片过滤：太小的图（< 100×100 px 或 < 5KB）自动跳过
- 扫描版 PDF 检测：文字为空但有图片内容 → 标记为需多模态处理

用法：
    python3 read_reference.py <文件路径>                    # 读取单个文件
    python3 read_reference.py <文件路径> --extract-images   # 读取并导出嵌入图片
    python3 read_reference.py <目录路径>                    # 列出目录中所有可读文件
    python3 read_reference.py <目录路径> --all              # 读取目录中所有文件内容
"""

import sys
import os
import hashlib
from pathlib import Path

# === 能力检测 ===

import subprocess

def _auto_install(package_name, import_name=None):
    """尝试导入，失败则自动 pip install 再导入"""
    import_name = import_name or package_name
    try:
        __import__(import_name)
        return True
    except ImportError:
        try:
            subprocess.check_call(
                [sys.executable, '-m', 'pip', 'install', '-q', package_name],
                stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
            )
            __import__(import_name)
            return True
        except Exception:
            return False

HAS_MARKITDOWN = _auto_install('markitdown')
HAS_DOCX = _auto_install('python-docx', 'docx')
HAS_PYPDF2 = _auto_install('PyPDF2', 'PyPDF2')
HAS_PYMUPDF = _auto_install('pymupdf', 'fitz')
HAS_PPTX = _auto_install('python-pptx', 'pptx')

if HAS_MARKITDOWN:
    from markitdown import MarkItDown
if HAS_DOCX:
    from docx import Document as DocxDocument
    from docx.opc.constants import RELATIONSHIP_TYPE as RT
if HAS_PYPDF2:
    import PyPDF2
if HAS_PYMUPDF:
    import fitz
if HAS_PPTX:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

# 支持的文件格式
TEXT_FORMATS = {'.md', '.txt', '.text', '.csv'}
DOCUMENT_FORMATS = {'.docx', '.pdf', '.pptx'}
IMAGE_FORMATS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'}
ALL_SUPPORTED = TEXT_FORMATS | DOCUMENT_FORMATS | IMAGE_FORMATS

# 图片过滤阈值
MIN_IMAGE_SIZE_KB = 5       # 小于 5KB 的图片视为装饰性图片
MIN_IMAGE_DIMENSION = 100   # 宽或高小于 100px 的图片视为装饰性图片


# ============================================================
# 图片导出工具
# ============================================================

def _ensure_image_dir(filepath):
    """为文件创建对应的图片导出目录"""
    filepath = Path(filepath)
    image_dir = filepath.parent / f"{filepath.stem}_images"
    image_dir.mkdir(exist_ok=True)
    return image_dir


def _is_decorative(image_bytes, width=None, height=None):
    """判断图片是否是装饰性的（太小不值得让模型看）"""
    # 文件大小过滤
    if len(image_bytes) < MIN_IMAGE_SIZE_KB * 1024:
        return True
    # 尺寸过滤（如果能拿到尺寸信息）
    if width and height:
        if width < MIN_IMAGE_DIMENSION or height < MIN_IMAGE_DIMENSION:
            return True
    return False


def _deduplicate_hash(image_bytes):
    """返回图片内容的哈希值，用于去重"""
    return hashlib.md5(image_bytes).hexdigest()[:12]


def extract_images_from_docx(filepath, image_dir):
    """从 Word 文档中提取嵌入图片"""
    if not HAS_DOCX:
        return []

    extracted = []
    seen_hashes = set()

    try:
        doc = DocxDocument(filepath)
        img_index = 0

        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    image_data = rel.target_part.blob
                    content_type = rel.target_part.content_type
                    ext = _content_type_to_ext(content_type)

                    if _is_decorative(image_data):
                        continue

                    img_hash = _deduplicate_hash(image_data)
                    if img_hash in seen_hashes:
                        continue
                    seen_hashes.add(img_hash)

                    img_index += 1
                    img_name = f"fig_{img_index}{ext}"
                    img_path = image_dir / img_name
                    img_path.write_bytes(image_data)

                    size_kb = len(image_data) / 1024
                    extracted.append({
                        'index': img_index,
                        'filename': img_name,
                        'path': str(img_path.resolve()),
                        'size_kb': size_kb,
                    })
                except Exception:
                    continue
    except Exception:
        pass

    return extracted


def extract_images_from_pdf(filepath, image_dir, max_pages=None):
    """从 PDF 中提取嵌入图片"""
    if not HAS_PYMUPDF:
        return [], False

    extracted = []
    seen_hashes = set()
    img_index = 0
    is_scanned = False

    try:
        doc = fitz.open(str(filepath))
        total_pages = len(doc)
        pages_to_check = min(total_pages, max_pages or total_pages)

        has_text = False
        has_full_page_images = False

        for page_num in range(pages_to_check):
            page = doc[page_num]
            page_text = page.get_text().strip()
            if page_text:
                has_text = True

            image_list = page.get_images(full=True)
            for img_info in image_list:
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    image_data = base_image["image"]
                    ext = f".{base_image['ext']}"
                    width = base_image.get("width", 0)
                    height = base_image.get("height", 0)

                    if _is_decorative(image_data, width, height):
                        continue

                    # 检测是否是整页扫描图（图片尺寸接近页面尺寸）
                    page_rect = page.rect
                    if width > page_rect.width * 0.8 and height > page_rect.height * 0.8:
                        has_full_page_images = True

                    img_hash = _deduplicate_hash(image_data)
                    if img_hash in seen_hashes:
                        continue
                    seen_hashes.add(img_hash)

                    img_index += 1
                    img_name = f"page{page_num + 1}_fig_{img_index}{ext}"
                    img_path = image_dir / img_name
                    img_path.write_bytes(image_data)

                    size_kb = len(image_data) / 1024
                    extracted.append({
                        'index': img_index,
                        'filename': img_name,
                        'path': str(img_path.resolve()),
                        'size_kb': size_kb,
                        'page': page_num + 1,
                        'width': width,
                        'height': height,
                    })
                except Exception:
                    continue

        doc.close()

        # 判断是否为扫描版 PDF
        if not has_text and has_full_page_images:
            is_scanned = True

    except Exception:
        pass

    return extracted, is_scanned


def extract_images_from_pptx(filepath, image_dir):
    """从 PPT 中提取嵌入图片（过滤装饰性图片）"""
    if not HAS_PPTX:
        return []

    extracted = []
    seen_hashes = set()
    img_index = 0

    try:
        prs = Presentation(str(filepath))

        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                # 跳过非图片形状
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue

                try:
                    image = shape.image
                    image_data = image.blob
                    ext = f".{image.content_type.split('/')[-1]}"
                    if ext == '.jpeg':
                        ext = '.jpg'

                    # 获取形状尺寸（EMU 转像素，1 inch = 914400 EMU, 96 DPI）
                    width_px = int(shape.width / 914400 * 96) if shape.width else 0
                    height_px = int(shape.height / 914400 * 96) if shape.height else 0

                    if _is_decorative(image_data, width_px, height_px):
                        continue

                    img_hash = _deduplicate_hash(image_data)
                    if img_hash in seen_hashes:
                        continue
                    seen_hashes.add(img_hash)

                    img_index += 1
                    img_name = f"slide{slide_num}_fig_{img_index}{ext}"
                    img_path = image_dir / img_name
                    img_path.write_bytes(image_data)

                    size_kb = len(image_data) / 1024
                    extracted.append({
                        'index': img_index,
                        'filename': img_name,
                        'path': str(img_path.resolve()),
                        'size_kb': size_kb,
                        'slide': slide_num,
                        'width': width_px,
                        'height': height_px,
                    })
                except Exception:
                    continue
    except Exception:
        pass

    return extracted


def _content_type_to_ext(content_type):
    """将 MIME 类型转换为文件扩展名"""
    mapping = {
        'image/png': '.png',
        'image/jpeg': '.jpg',
        'image/gif': '.gif',
        'image/bmp': '.bmp',
        'image/webp': '.webp',
        'image/tiff': '.tiff',
        'image/x-emf': '.emf',
        'image/x-wmf': '.wmf',
    }
    return mapping.get(content_type, '.png')


# ============================================================
# 文本读取（原有功能，小幅调整）
# ============================================================

def read_text_file(filepath):
    """读取纯文本文件（尝试多种编码）"""
    for enc in ['utf-8', 'gbk', 'gb2312', 'gb18030', 'latin-1']:
        try:
            return Path(filepath).read_text(encoding=enc)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return f"[错误] 无法解码文件：{filepath}"


def read_docx(filepath):
    """读取 Word 文档的文字内容"""
    if HAS_MARKITDOWN:
        try:
            md = MarkItDown()
            result = md.convert_local(str(filepath))
            text = getattr(result, 'text_content', None) or getattr(result, 'markdown', str(result))
            if text and len(text.strip()) > 10:
                return text
        except Exception:
            pass

    if HAS_DOCX:
        try:
            doc = DocxDocument(filepath)
            parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                    if row_text.strip('| '):
                        parts.append(row_text)
            if parts:
                return '\n\n'.join(parts)
        except Exception:
            pass

    return f"[跳过] 无法读取 Word 文件 {Path(filepath).name}（解析库不可用）。"


def read_pdf(filepath):
    """读取 PDF 文件的文字内容"""
    if HAS_MARKITDOWN:
        try:
            md = MarkItDown()
            result = md.convert_local(str(filepath))
            text = getattr(result, 'text_content', None) or getattr(result, 'markdown', str(result))
            if text and len(text.strip()) > 10:
                return text
        except Exception:
            pass

    if HAS_PYPDF2:
        try:
            parts = []
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        parts.append(page_text)
            if parts:
                return '\n\n'.join(parts)
        except Exception:
            pass

    return ""


def read_pptx(filepath):
    """读取 PPT 文件的文字内容"""
    if HAS_MARKITDOWN:
        try:
            md = MarkItDown()
            result = md.convert_local(str(filepath))
            text = getattr(result, 'text_content', None) or getattr(result, 'markdown', str(result))
            if text and len(text.strip()) > 10:
                return text
        except Exception:
            pass

    return f"[跳过] 无法读取 PPT 文件 {Path(filepath).name}（解析库不可用）。"


# ============================================================
# 统一入口
# ============================================================

def read_file(filepath, extract_images=False):
    """
    读取单个文件，返回文本内容。
    如果 extract_images=True，同时导出嵌入图片并在输出中标记索引。
    """
    filepath = Path(filepath)
    suffix = filepath.suffix.lower()

    if not filepath.exists():
        return f"[错误] 文件不存在：{filepath}"

    # 纯文本
    if suffix in TEXT_FORMATS:
        return read_text_file(filepath)

    # 独立图片文件
    if suffix in IMAGE_FORMATS:
        size_kb = filepath.stat().st_size / 1024
        return (
            f"[图片] {filepath.name} ({size_kb:.0f} KB)\n"
            f"这是一张独立的图片文件，需要多模态模型用 Read 工具查看。\n"
            f"图片路径：{filepath.resolve()}"
        )

    # === Word / PDF / PPT：文字提取 + 可选的图片导出 ===

    text = ""
    images = []
    is_scanned = False

    if suffix == '.docx':
        text = read_docx(filepath)
        if extract_images:
            image_dir = _ensure_image_dir(filepath)
            images = extract_images_from_docx(filepath, image_dir)

    elif suffix == '.pdf':
        text = read_pdf(filepath)
        if extract_images:
            image_dir = _ensure_image_dir(filepath)
            images, is_scanned = extract_images_from_pdf(filepath, image_dir)

    elif suffix == '.pptx':
        text = read_pptx(filepath)
        if extract_images:
            image_dir = _ensure_image_dir(filepath)
            images = extract_images_from_pptx(filepath, image_dir)

    else:
        return f"[提示] 不支持的格式：{suffix}"

    # === 组装输出 ===

    output_parts = []

    # 扫描版 PDF 提示
    if is_scanned:
        output_parts.append(
            f"[扫描版 PDF] {filepath.name} 是扫描版文档（纯图像，无可提取的文字）。\n"
            f"需要多模态模型用 Read 工具按页查看：Read(\"{filepath.resolve()}\", pages=\"1-5\")\n"
            f"共 {len(images)} 张页面图片已导出到 {image_dir}/"
        )
    else:
        # 文字内容
        if text and text.strip():
            output_parts.append(text)
        else:
            output_parts.append(f"[提示] 未能从 {filepath.name} 中提取到文字内容。")

    # 图片索引
    if images and not is_scanned:
        output_parts.append("")
        output_parts.append(f"--- 嵌入图片（共 {len(images)} 张，已导出） ---")
        output_parts.append("以下图片从文档中提取，需要多模态模型用 Read 工具查看并理解内容：")
        output_parts.append("")
        for img in images:
            location = ""
            if 'page' in img:
                location = f"第{img['page']}页 · "
            elif 'slide' in img:
                location = f"第{img['slide']}张幻灯片 · "
            size_info = f"{img['size_kb']:.0f}KB"
            if 'width' in img and 'height' in img:
                size_info += f" · {img['width']}×{img['height']}px"
            output_parts.append(f"  [图片 {img['index']}] {location}{size_info}")
            output_parts.append(f"    路径：{img['path']}")

    return '\n'.join(output_parts)


def list_directory(dirpath):
    """列出目录中所有可读的参考文件"""
    dirpath = Path(dirpath)
    if not dirpath.is_dir():
        print(f"[错误] 不是有效目录：{dirpath}")
        return

    files = []
    for f in sorted(dirpath.rglob('*')):
        if f.is_file() and f.suffix.lower() in ALL_SUPPORTED and not f.name.startswith('.'):
            rel_path = f.relative_to(dirpath)
            size_kb = f.stat().st_size / 1024
            suffix = f.suffix.lower()

            if suffix in IMAGE_FORMATS:
                icon = "🖼️"
                note = "（需多模态模型）"
            elif suffix in DOCUMENT_FORMATS:
                icon = "📄"
                note = ""
            else:
                icon = "📝"
                note = ""

            files.append(f"  {icon} {rel_path}  ({size_kb:.0f} KB) {note}")

    # 输出环境信息
    print(f"参考资料目录：{dirpath}")
    can_read = []
    if HAS_MARKITDOWN or HAS_DOCX:
        can_read.append('Word')
    if HAS_MARKITDOWN or HAS_PYPDF2 or HAS_PYMUPDF:
        can_read.append('PDF')
    if HAS_MARKITDOWN or HAS_PPTX:
        can_read.append('PPT')
    can_extract_img = []
    if HAS_DOCX:
        can_extract_img.append('Word')
    if HAS_PYMUPDF:
        can_extract_img.append('PDF')
    if HAS_PPTX:
        can_extract_img.append('PPT')

    print(f"文字提取：{', '.join(can_read) if can_read else '不可用'}")
    print(f"图片导出：{', '.join(can_extract_img) if can_extract_img else '不可用'}（--extract-images 启用）")
    print()

    if files:
        print(f"共 {len(files)} 个文件：")
        for f in files:
            print(f)
    else:
        print("目录中没有找到可读的参考文件。")
        print(f"支持的格式：{', '.join(sorted(ALL_SUPPORTED))}")


def read_all_in_directory(dirpath, extract_images=False):
    """读取目录中所有文件的内容"""
    dirpath = Path(dirpath)
    if not dirpath.is_dir():
        print(f"[错误] 不是有效目录：{dirpath}")
        return

    files = sorted(f for f in dirpath.rglob('*')
                   if f.is_file()
                   and f.suffix.lower() in (TEXT_FORMATS | DOCUMENT_FORMATS)
                   and not f.name.startswith('.'))

    image_files = sorted(f for f in dirpath.rglob('*')
                         if f.is_file()
                         and f.suffix.lower() in IMAGE_FORMATS
                         and not f.name.startswith('.'))

    if not files and not image_files:
        print("目录中没有可读的文件。")
        return

    for f in files:
        rel_path = f.relative_to(dirpath)
        print(f"\n{'='*60}")
        print(f"📄 {rel_path}")
        print(f"{'='*60}\n")
        content = read_file(f, extract_images=extract_images)
        if len(content) > 8000:
            print(content[:8000])
            print(f"\n... [截断，完整 {len(content)} 字]")
        else:
            print(content)

    if image_files:
        print(f"\n{'='*60}")
        print(f"🖼️ 独立图片文件（共 {len(image_files)} 张，需 Read 工具查看）：")
        print(f"{'='*60}\n")
        for f in image_files:
            rel_path = f.relative_to(dirpath)
            size_kb = f.stat().st_size / 1024
            print(f"  {rel_path}  ({size_kb:.0f} KB) → {f.resolve()}")


def main():
    if len(sys.argv) < 2:
        print("读取参考资料工具（支持文字提取 + 图片导出）")
        print()
        print("用法：")
        print("  python3 read_reference.py <文件>                    读取单个文件（仅文字）")
        print("  python3 read_reference.py <文件> --extract-images   读取并导出嵌入图片")
        print("  python3 read_reference.py <目录>                    列出可读文件")
        print("  python3 read_reference.py <目录> --all              读取所有文件")
        print("  python3 read_reference.py <目录> --all --extract-images  读取所有文件并导出图片")
        print()
        print("图片导出说明：")
        print("  --extract-images 会将 Word/PDF/PPT 中的嵌入图片导出为独立文件，")
        print("  并在输出中标记索引。导出的图片需要多模态模型用 Read 工具查看。")
        print(f"  装饰性图片（< {MIN_IMAGE_SIZE_KB}KB 或 < {MIN_IMAGE_DIMENSION}×{MIN_IMAGE_DIMENSION}px）自动过滤。")
        sys.exit(0)

    target = Path(sys.argv[1])
    read_all = '--all' in sys.argv
    extract_images = '--extract-images' in sys.argv

    if target.is_dir():
        if read_all:
            read_all_in_directory(target, extract_images=extract_images)
        else:
            list_directory(target)
    elif target.is_file():
        print(read_file(target, extract_images=extract_images))
    else:
        print(f"[错误] 路径不存在：{target}")
        sys.exit(1)


if __name__ == '__main__':
    main()
